import pandas as pd
import os
import re
from pathlib import Path
import subprocess
import sys
import tempfile
import time
from pptx import Presentation
import shutil

def sanitize_filename(filename):
    """Remove caracteres inválidos para nomes de arquivos"""
    return re.sub(r'[<>:"/\\|?*]', '_', str(filename))

def generate_pdf_with_libreoffice(template_path, output_pdf, replacements):
    """
    Gera PDF diretamente do template usando LibreOffice
    :param template_path: Caminho para o template PPTX
    :param output_pdf: Caminho de saída para o PDF
    :param replacements: Dicionário com {placeholder: valor}
    """
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, "temp_certificado.pptx")
    temp_pdf = os.path.join(temp_dir, "temp_certificado.pdf")
    
    try:
        # 1. Copiar o template para arquivo temporário
        shutil.copyfile(template_path, temp_pptx)
        
        # 2. Modificar o template temporário
        prs = Presentation(temp_pptx)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for ph, value in replacements.items():
                                if ph in run.text:
                                    run.text = run.text.replace(ph, value)
        
        prs.save(temp_pptx)
        
        # 3. Converter para PDF usando LibreOffice
        soffice_path = None
        possible_paths = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            r'C:\Arquivos de Programas\LibreOffice\program\soffice.exe'
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                soffice_path = path
                break
        
        if not soffice_path:
            raise FileNotFoundError("LibreOffice não encontrado. Verifique a instalação.")
        
        cmd = [
            soffice_path,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', temp_dir,
            temp_pptx
        ]
        
        # Executar a conversão
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, 
                              text=True, timeout=60, shell=True)
        
        # 4. Verificar se a conversão foi bem-sucedida
        if not os.path.exists(temp_pdf):
            print("Erro: Arquivo PDF não foi gerado")
            print("Saída do LibreOffice:")
            print(result.stdout)
            print(result.stderr)
            return False
        
        # 5. Mover para o destino final
        if os.path.exists(output_pdf):
            os.remove(output_pdf)
        shutil.move(temp_pdf, output_pdf)
        
        return True
        
    except Exception as e:
        print(f"Erro durante a geração do PDF: {str(e)}")
        return False
    finally:
        # Limpeza dos arquivos temporários
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except:
            pass

def gerar_certificados(template_path, csv_path, output_folder):
    # Verificações iniciais
    print(f"\nVerificando arquivos...")
    print(f"Template existe? {os.path.exists(template_path)}")
    print(f"CSV existe? {os.path.exists(csv_path)}")
    
    # Criar pasta de saída
    Path(output_folder).mkdir(parents=True, exist_ok=True)
    
    # Ler CSV
    try:
        df = pd.read_csv(csv_path)
        print(f"\nTotal de registros encontrados: {len(df)}")
        print("\nAmostra dos dados:")
        print(df.head())
    except Exception as e:
        print(f"\nERRO ao ler CSV: {e}")
        return

    # Contadores
    sucessos = 0
    erros = 0
    start_time = time.time()

    # Processar cada registro
    for index, row in df.iterrows():
        try:
            nome = str(row['nome']).strip()
            numero = str(row['numero']).strip()
            
            print(f"\nProcessando: {nome} (Certificado {numero})")
            
            # Sanitizar valores e formatar número
            nome_safe = sanitize_filename(nome)
            numero_safe = sanitize_filename(numero).replace('/', '_')
            
            # Criar nome de arquivo seguro
            output_filename = f"Certificado_{numero_safe}_{nome_safe[:50]}.pdf"
            output_path = os.path.join(output_folder, output_filename)
            
            # Gerar PDF diretamente
            replacements = {
                '{NOME}': nome,
                '{NUMERO}': numero
            }
            
            if generate_pdf_with_libreoffice(template_path, output_path, replacements):
                print(f"✓ PDF gerado com sucesso: {os.path.basename(output_path)}")
                sucessos += 1
            else:
                print(f"✗ Falha na geração do PDF para {nome}")
                erros += 1
            
        except Exception as e:
            print(f"ERRO ao processar linha {index + 1}: {e}")
            erros += 1
    
    # Relatório final
    total_time = (time.time() - start_time) / 60
    print("\n" + "="*50)
    print("RELATÓRIO FINAL")
    print(f"Total de certificados processados: {len(df)}")
    print(f"PDFs gerados com sucesso: {sucessos}")
    print(f"Erros encontrados: {erros}")
    if sucessos > 0:
        print(f"\nArquivos gerados em: {os.path.abspath(output_folder)}")
    print("="*50)

if __name__ == "__main__":
        # CONFIGURAÇÕES - AJUSTE ESTES CAMINHOS!
    TEMPLATE_PATH = r"./base/certificado_template.pptx"
    CSV_PATH = r"./base/dados_participantes.csv"
    OUTPUT_FOLDER = r"./certificados_gerados"
    
    print("Iniciando processo de geração de certificados em PDF...")
    gerar_certificados(TEMPLATE_PATH, CSV_PATH, OUTPUT_FOLDER)
    
    # Manter o prompt aberto
    input("\nPressione ENTER para fechar esta janela...")
    
